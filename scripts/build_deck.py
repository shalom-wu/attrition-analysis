"""
Build reports/strategy_deck.pptx from the analysis results.

Mirrors reports/strategy_deck.md: 9 slides (title + 8 content). Charts are
native PowerPoint charts built from the pipeline outputs; the SHAP summary
is embedded as an image (no native chart type exists for it).

Palette ("Midnight Executive"): navy 1E2761 dominant, ice blue CADCFC,
white, coral E14B4B as the risk accent. Fonts: Cambria heads / Calibri body.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "figures"

NAVY = RGBColor.from_string("1E2761")
ICE = RGBColor.from_string("CADCFC")
WHITE = RGBColor.from_string("FFFFFF")
CORAL = RGBColor.from_string("E14B4B")
DARKTXT = RGBColor.from_string("232946")
MUTED = RGBColor.from_string("5A6178")
TINT = RGBColor.from_string("EEF3FC")     # light navy tint for cards
GREEN = RGBColor.from_string("2E7D32")
GRID = RGBColor.from_string("E2E8F0")

HEAD = "Cambria"
BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------
def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf

def para(tf, text, size=14, color=DARKTXT, bold=False, italic=False,
         font=BODY, align=PP_ALIGN.LEFT, space_after=6, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    f = r.font
    f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
    f.color.rgb = color
    return p

def runs(tf, parts, size=14, align=PP_ALIGN.LEFT, space_after=6, first=False):
    """parts = [(text, {overrides}), ...] on one paragraph."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    for text, ov in parts:
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", False)
        f.italic = ov.get("italic", False)
        f.name = ov.get("font", BODY)
        f.color.rgb = ov.get("color", DARKTXT)
    return p

def title(slide, text, color=NAVY, y=0.42, size=30):
    tf = box(slide, 0.6, y, 12.1, 0.85)
    para(tf, text, size=size, color=color, bold=True, font=HEAD, first=True)

def card(slide, x, y, w, h, fill=TINT):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.adjustments[0] = 0.06
    sh.shadow.inherit = False
    return sh

def stat(slide, x, y, w, number, label, num_color=WHITE, lab_color=ICE,
         num_size=40, lab_size=11):
    tf = box(slide, x, y, w, 1.3)
    para(tf, number, size=num_size, color=num_color, bold=True, font=HEAD,
         align=PP_ALIGN.CENTER, space_after=2, first=True)
    para(tf, label, size=lab_size, color=lab_color, align=PP_ALIGN.CENTER)

def style_chart(chart, colors, number_format="0.0", label_size=10,
                show_labels=True, gap=80):
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = gap
    for i, s in enumerate(plot.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = colors[i % len(colors)]
    if show_labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = number_format
        dl.number_format_is_linked = False
        dl.font.size = Pt(label_size)
        dl.font.bold = True
        dl.font.color.rgb = DARKTXT
        dl.font.name = BODY
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    ca = chart.category_axis
    va = chart.value_axis
    ca.tick_labels.font.size = Pt(10)
    ca.tick_labels.font.name = BODY
    ca.tick_labels.font.color.rgb = MUTED
    ca.has_major_gridlines = False
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = GRID
    va.major_gridlines.format.line.width = Pt(0.5)
    va.tick_labels.font.size = Pt(9)
    va.tick_labels.font.name = BODY
    va.tick_labels.font.color.rgb = MUTED
    va.format.line.fill.background()
    ca.format.line.color.rgb = GRID

def chart_title(chart, text):
    chart.has_title = True
    chart.chart_title.text_frame.text = text
    r = chart.chart_title.text_frame.paragraphs[0].runs[0]
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.name = BODY
    r.font.color.rgb = DARKTXT

# ===========================================================================
# Slide 1 — Title (navy)
# ===========================================================================
s = add_slide(NAVY)
tf = box(s, 0.9, 1.15, 11.5, 0.4)
para(tf, "PEOPLE ANALYTICS  ·  JULY 2026", size=13, color=ICE, bold=True, first=True)
tf = box(s, 0.9, 1.75, 11.5, 1.9)
para(tf, "Employee Attrition: Cost, Drivers,", size=42, color=WHITE, bold=True,
     font=HEAD, space_after=0, first=True)
para(tf, "and a Retention Plan", size=42, color=WHITE, bold=True, font=HEAD)
tf = box(s, 0.9, 3.55, 11.5, 0.5)
para(tf, "Analysis of 1,470 employees across Engineering, Sales, and HR — "
         "descriptive analysis, predictive model, and costed interventions.",
     size=15, color=ICE, first=True)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(4.7),
                          Inches(11.5), Emu(9525))
line.fill.solid(); line.fill.fore_color.rgb = RGBColor.from_string("3A4580")
line.line.fill.background(); line.shadow.inherit = False
stat(s, 0.9, 5.15, 3.6, "16.1%", "ANNUAL ATTRITION — 237 OF 1,470 EMPLOYEES")
stat(s, 4.85, 5.15, 3.6, "$8.67M", "ANNUAL REPLACEMENT COST — 7.6% OF PAYROLL",
     num_color=CORAL)
stat(s, 8.8, 5.15, 3.6, "7.1 mo", "OF SALARY LOST PER DEPARTURE (AVG $36.6K)")

# ===========================================================================
# Slide 2 — The problem
# ===========================================================================
s = add_slide()
title(s, "Attrition costs us ~$8.7M a year — and it's concentrated")
tf = box(s, 0.6, 1.5, 5.6, 5.4)
items = [
    ("16.1% annual attrition", " sits at the high end of the 13–17% typical "
     "range for mid-size tech/SaaS companies."),
    ("$36.6K per departure", " — recruiting, onboarding, vacancy, and ramp-up; "
     "≈7.1 months of salary, conservative vs. SHRM (6–9 mo) and Gallup "
     "(0.5–2× salary) benchmarks."),
    ("Sales bleeds fastest", " at 20.6% attrition ($4.25M/yr); Engineering "
     "loses the most absolute dollars after Sales ($4.10M/yr at 13.8%)."),
    ("Every dollar is traceable", " to an explicit assumption — soft costs "
     "(knowledge loss, morale, customer damage) are deliberately excluded, "
     "so the true number is higher."),
]
first = True
for lead, rest in items:
    runs(tf, [(lead, {"bold": True, "color": NAVY}), (rest, {"color": DARKTXT})],
         size=14.5, space_after=14, first=first)
    first = False

data = CategoryChartData()
data.categories = ["Sales", "Engineering", "HR"]
data.add_series("Annual cost ($M)", (4.25, 4.10, 0.32))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.6), Inches(1.5),
                        Inches(6.1), Inches(4.1), data)
style_chart(gf.chart, [NAVY], number_format='"$"0.00"M"')
chart_title(gf.chart, "Annual replacement cost by department")
c = card(s, 6.6, 5.85, 6.1, 1.05)
tf = c.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.2)
runs(tf, [("Attrition rate by department:  ", {"bold": True, "color": NAVY}),
          ("Sales 20.6%   ·   HR 19.0%   ·   Engineering 13.8%   "
           "(company avg 16.1%)", {"color": DARKTXT})], size=13, first=True)

# ===========================================================================
# Slide 3 — Where attrition concentrates
# ===========================================================================
s = add_slide()
title(s, "It's not random: overtime and early tenure do the damage")
data = CategoryChartData()
data.categories = ["L1", "L2", "L3", "L4", "L5"]
data.add_series("No overtime", (15.8, 6.7, 12.3, 2.7, 3.9))
data.add_series("Overtime", (52.6, 17.8, 20.6, 9.1, 16.7))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.5),
                        Inches(6.6), Inches(4.4), data)
ch = gf.chart
style_chart(ch, [NAVY, CORAL], number_format='0.0"%"', gap=60)
chart_title(ch, "Attrition rate (%) by job level × overtime")
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11)
ch.legend.font.name = BODY

c = card(s, 7.6, 1.5, 5.1, 2.0, fill=RGBColor.from_string("FDEDED"))
tf = c.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.25)
tf.margin_top = Inches(0.18)
para(tf, "52.6%", size=38, color=CORAL, bold=True, font=HEAD, space_after=2, first=True)
para(tf, "attrition among junior (L1) employees working overtime — a third of "
         "all departures from just 11% of headcount.", size=13, color=DARKTXT)

tf = box(s, 7.75, 3.85, 5.0, 3.3)
rows = [
    ("34.9%", "first-year employees (4.3× the 10-yr+ rate)"),
    ("30.5%", "overtime workers overall (vs 10.4% without)"),
    ("29.3%", "bottom income quartile (2.8× top quartile)"),
    ("24.9%", "frequent business travelers (3.1× non-travelers)"),
    ("66.7%", "sales representatives working overtime"),
]
first = True
for num, lab in rows:
    runs(tf, [(num + "   ", {"bold": True, "color": NAVY, "size": 17, "font": HEAD}),
              (lab, {"color": MUTED, "size": 12.5})],
         space_after=10, first=first)
    first = False

# ===========================================================================
# Slide 4 — Model
# ===========================================================================
s = add_slide()
title(s, "The model confirms the drivers — and adds a hidden one", size=28)
tf = box(s, 0.6, 1.55, 5.9, 4.3)
drivers = [
    ("1. Overtime", " — strongest predictor; ~2× the odds of leaving, all else equal."),
    ("2. Low pay / junior level / short tenure", " — a correlated cluster; top "
     "SHAP features after overtime (income, stock options, age)."),
    ("3. Frequent business travel", " — odds ratio ≈1.7; underappreciated."),
    ("4. Promotion stagnation", " — years since last promotion raises risk "
     "(OR ≈1.6/SD) once tenure is controlled for. The naive segment cut hides "
     "this; the regression surfaces it."),
    ("5. Low satisfaction, long commute, single", " — secondary but consistent."),
]
first = True
for lead, rest in drivers:
    runs(tf, [(lead, {"bold": True, "color": NAVY}), (rest, {"color": DARKTXT})],
         size=13.5, space_after=17, first=first)
    first = False

c = card(s, 0.6, 5.85, 5.9, 1.15)
tf = c.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.2)
runs(tf, [("Model quality:  ", {"bold": True, "color": NAVY}),
          ("ROC-AUC 0.81 (logistic; XGBoost 0.79). Catches 71% of true "
           "leavers while flagging 33% of employees — for prioritizing "
           "manager check-ins, not automated decisions.", {"color": DARKTXT})],
     size=12.5, first=True)

# SHAP summary image: 1024x966 px → keep aspect ratio
img_h = 5.55
img_w = img_h * (1024 / 966)
s.shapes.add_picture(str(FIG / "shap_summary.png"),
                     Inches(13.333 - 0.5 - img_w), Inches(1.45),
                     Inches(img_w), Inches(img_h))

# ===========================================================================
# Slide 5 — Cost of inaction
# ===========================================================================
s = add_slide()
title(s, "Cost of inaction: ~$26M over the next three years")
tbl_rows = [
    ("Population", "Leavers / yr", "Annual cost", "Share of cost"),
    ("Overtime employees (n=416)", "127", "$4.55M", "52.5%"),
    ("First-year employees (n=215)", "75", "$1.68M", "19.4%"),
    ("Junior + overtime core (n=156)", "82", "$1.42M", "16.3%"),
    ("Everyone else", "—", "balance", "—"),
]
gfx = s.shapes.add_table(5, 4, Inches(0.6), Inches(1.6), Inches(7.4), Inches(3.1))
tbl = gfx.table
tbl.columns[0].width = Inches(3.4)
tbl.columns[1].width = Inches(1.3)
tbl.columns[2].width = Inches(1.4)
tbl.columns[3].width = Inches(1.3)
for ri, row in enumerate(tbl_rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        f = p.runs[0].font
        f.name = BODY
        f.size = Pt(12.5)
        if ri == 0:
            f.bold = True; f.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        else:
            f.color.rgb = DARKTXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else TINT

tf = box(s, 0.6, 5.15, 7.4, 1.9)
para(tf, "Populations overlap (a first-year junior on overtime appears in all "
         "three rows) — shares don't sum to 100%.", size=11.5, color=MUTED,
     italic=True, space_after=10, first=True)
runs(tf, [("What's not in the number:  ", {"bold": True, "color": NAVY}),
          ("lost domain knowledge, customer relationship damage in Sales, "
           "team morale drag, manager time. All push the true cost higher.",
           {"color": DARKTXT})], size=13)

c = card(s, 8.5, 1.6, 4.2, 4.0, fill=NAVY)
tf = c.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
tf.margin_top = Inches(0.35)
para(tf, "$8.67M", size=44, color=WHITE, bold=True, font=HEAD,
     align=PP_ALIGN.CENTER, space_after=2, first=True)
para(tf, "per year at the current run-rate", size=12, color=ICE,
     align=PP_ALIGN.CENTER, space_after=18)
para(tf, "$26M", size=44, color=CORAL, bold=True, font=HEAD,
     align=PP_ALIGN.CENTER, space_after=2)
para(tf, "over three years, before headcount growth", size=12, color=ICE,
     align=PP_ALIGN.CENTER)

# ===========================================================================
# Slide 6 — Options
# ===========================================================================
s = add_slide()
title(s, "Three intervention options")
options = [
    ("A — Targeted overtime reduction",
     "On-call restructuring, workload rebalancing, comp-time, and a ~4-hire "
     "backfill buffer, prioritizing the 156 junior+overtime core.",
     [("Effect (assumed):", " −25% relative attrition in the OT group → 32 fewer exits, $1.13M/yr saved"),
      ("Cost:", " ~$300K/yr")],
     "Net ≈ +$830K/yr  ·  ~3.8× gross ROI", GREEN,
     "Tradeoff: manager behavior change; some delivery-speed sacrifice."),
    ("B — Broad comp adjustment",
     "Raise the 369 bottom-quartile earners by 10%.",
     [("Effect (assumed):", " −25% relative attrition → 27 fewer exits, but junior replacements are cheap: $0.40M/yr saved"),
      ("Cost:", " ~$1.04M/yr, recurring")],
     "Net ≈ −$640K/yr  ·  does not pay for itself", CORAL,
     "Tradeoff: strongest signal of investment; hardest to reverse. Reserve "
     "targeted corrections for proven market gaps (needs comp benchmark)."),
    ("C — First-year experience + promotion velocity",
     "Structured 90-day onboarding, mentor program, defined L1→L2 path with "
     "12–18 month checkpoints.",
     [("Effect (assumed):", " −20% first-year attrition → 15 fewer exits, $0.34M/yr saved + promotion-velocity upside"),
      ("Cost:", " ~$120K/yr")],
     "Net ≈ +$215K/yr  ·  ~2.8× gross ROI", GREEN,
     "Tradeoff: cheapest, lowest-risk; addresses only ~20% of the cost pool."),
]
xs = [0.6, 4.93, 9.26]
for (name, desc, lines, net, net_color, trade), x in zip(options, xs):
    card(s, x, 1.75, 4.07, 4.1)
    tf = box(s, x + 0.28, 2.05, 3.51, 3.6)
    para(tf, name, size=15.5, color=NAVY, bold=True, font=HEAD,
         space_after=8, first=True)
    para(tf, desc, size=11.5, color=DARKTXT, space_after=10)
    for lead, rest in lines:
        runs(tf, [(lead + " ", {"bold": True, "color": NAVY, "size": 11.5}),
                  (rest, {"color": DARKTXT, "size": 11.5})], space_after=8)
    para(tf, net, size=13, color=net_color, bold=True, space_after=8)
    para(tf, trade, size=10.5, color=MUTED, italic=True)

# ===========================================================================
# Slide 7 — Recommendation (navy)
# ===========================================================================
s = add_slide(NAVY)
title(s, "Recommendation: do A + C together — piloted, model-assisted",
      color=WHITE, size=28)
tf = box(s, 0.6, 1.55, 6.4, 2.2)
para(tf, "+$0.8M/yr net", size=48, color=WHITE, bold=True, font=HEAD,
     space_after=2, first=True)
para(tf, "≈ $1.25M/yr saved (overlap-discounted) for ~$420K/yr program cost "
         "— roughly 3× ROI", size=14, color=ICE)
tf = box(s, 0.6, 3.75, 6.4, 3.3)
recs = [
    ("~47 fewer departures a year", " — company attrition moves from 16.1% to "
     "~13%, from worst quartile to mid-band for tech."),
    ("Model sequences the rollout", " — risk scores pick which teams get "
     "on-call restructuring and enhanced onboarding first. Support only, "
     "never punitive."),
    ("Pilot before scaling", " — one quarter in the 2 highest-attrition Sales "
     "pods + 1 Engineering group, with a control group to validate the "
     "assumed effect sizes."),
    ("Defer B", " — revisit targeted comp corrections once market benchmark "
     "data is in."),
]
first = True
for lead, rest in recs:
    runs(tf, [(lead, {"bold": True, "color": WHITE}), (rest, {"color": ICE})],
         size=13, space_after=11, first=first)
    first = False

c = card(s, 7.6, 1.55, 5.1, 4.6, fill=RGBColor.from_string("2A3573"))
tf = c.text_frame; tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.TOP
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
tf.margin_top = Inches(0.3)
para(tf, "Stated assumptions", size=15, color=WHITE, bold=True, font=HEAD,
     space_after=10, first=True)
assumps = [
    "Effect sizes (−25% / −20% relative) are industry-plausible but unproven "
    "here — the pilot exists to test them.",
    "Replacement costs are conservative; soft costs excluded.",
    "Salary figures are the dataset's — if real payroll is higher, savings "
    "scale linearly.",
    "Savings discounted ~15% for population overlap between A and C.",
]
for a in assumps:
    para(tf, a, size=12, color=ICE, space_after=10)

# ===========================================================================
# Slide 8 — Roadmap
# ===========================================================================
s = add_slide()
title(s, "12-month roadmap")
road = [
    ("Q1", "Overtime audit + pilot in 3 highest-risk teams; structured "
     "onboarding for all new hires", "HRBP + Eng/Sales leads",
     "Pilot vs. control attrition; onboarding NPS"),
    ("Q2", "Read pilot; scale Option A if ≥15% relative reduction observed; "
     "define L1→L2 promotion rubric", "CHRO",
     "Effect size confirmed; rubric shipped"),
    ("Q3", "Full rollout; model-scored quarterly retention check-ins for the "
     "top-risk decile", "People Analytics",
     "71%+ of eventual leavers had a check-in"),
    ("Q4", "Re-measure: target ≤14% trailing attrition; decide targeted comp "
     "corrections with market data", "CHRO + Finance",
     "Attrition rate; $ saved vs. model"),
]
gfx = s.shapes.add_table(5, 4, Inches(0.6), Inches(1.6), Inches(12.13), Inches(4.6))
tbl = gfx.table
tbl.columns[0].width = Inches(0.85)
tbl.columns[1].width = Inches(5.6)
tbl.columns[2].width = Inches(2.4)
tbl.columns[3].width = Inches(3.28)
tbl.rows[0].height = Inches(0.45)
for ri in range(1, 5):
    tbl.rows[ri].height = Inches(1.0)
headers = ("", "Action", "Owner", "Success metric")
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.text = h or "Qtr"
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]
    f = p.runs[0].font
    f.bold = True; f.color.rgb = WHITE; f.size = Pt(13); f.name = BODY
for ri, (q, action, owner, metric) in enumerate(road, start=1):
    for ci, val in enumerate((q, action, owner, metric)):
        cell = tbl.cell(ri, ci)
        cell.text = val
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if ri % 2 else TINT
        p = cell.text_frame.paragraphs[0]
        f = p.runs[0].font
        f.size = Pt(12); f.name = BODY
        f.color.rgb = NAVY if ci == 0 else DARKTXT
        f.bold = ci == 0

# ===========================================================================
# Slide 9 — Appendix
# ===========================================================================
s = add_slide()
title(s, "Appendix: methodology & caveats", size=26)
tf = box(s, 0.6, 1.4, 5.9, 5.7)
para(tf, "Data", size=14, color=NAVY, bold=True, font=HEAD, space_after=4, first=True)
para(tf, "IBM HR Analytics dataset (1,470 employees, 35 fields, no missing "
         "values; 3 constant columns dropped). 'R&D' relabelled 'Engineering' "
         "for narrative. Synthetic and cross-sectional — a single snapshot, no "
         "dates — so attrition is treated as an annualized rate and no "
         "time-based validation is possible.", size=11.5, color=DARKTXT,
     space_after=12)
para(tf, "Cost model (per departure)", size=14, color=NAVY, bold=True,
     font=HEAD, space_after=4)
para(tf, "Recruiting 20% of annual salary + onboarding 10% + vacancy (1.5–4 "
         "months × monthly salary, by level) + ramp-up (3–6 months at 50% "
         "productivity). Yields 55% of salary (junior) to 88% (exec); average "
         "$36.6K ≈ 7.1 months. Anchors: SHRM cost-per-hire and time-to-fill "
         "studies; SHRM 6–9-month rule of thumb; Gallup 0.5–2× salary.",
     size=11.5, color=DARKTXT, space_after=12)
para(tf, "Models", size=14, color=NAVY, bold=True, font=HEAD, space_after=4)
para(tf, "75/25 stratified split, 5-fold CV; class imbalance handled with "
         "class weights (accuracy deliberately not headlined). Logistic "
         "regression: test ROC-AUC 0.810, PR-AUC 0.562 (baseline 0.16), CV "
         "0.842 ± 0.044. XGBoost: 0.785 / 0.522. The simpler model wins on "
         "this small dataset; both reported.", size=11.5, color=DARKTXT)

tf = box(s, 7.1, 1.4, 5.6, 5.7)
para(tf, "Caveats", size=14, color=NAVY, bold=True, font=HEAD, space_after=4,
     first=True)
caveats = [
    "Correlation ≠ causation — every intervention effect size is an "
    "assumption to be piloted, not a model output.",
    "237 positive cases → wide confidence intervals on per-role estimates.",
    "No voluntary/involuntary split — some 'attrition' may be terminations.",
    "No external comp benchmarks — can't distinguish 'paid low' from 'paid "
    "below market'.",
    "Ethics: risk scores prioritize support (check-ins, workload relief), "
    "never punitive action; not visible in performance contexts.",
]
for cv in caveats:
    para(tf, "–  " + cv, size=11.5, color=DARKTXT, space_after=8)
para(tf, "Open inputs needed from leadership", size=14, color=NAVY, bold=True,
     font=HEAD, space_after=4)
para(tf, "Real payroll scaling  ·  market comp benchmarks  ·  "
         "voluntary/involuntary attrition split  ·  actual time-to-fill by "
         "role.", size=11.5, color=DARKTXT)

out = ROOT / "reports" / "strategy_deck.pptx"
prs.save(out)
print(f"Wrote {out}")
